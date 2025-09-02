import sys, os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pathlib

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm
import requests
from io import BytesIO


def set_slide_background_picture(slide, prs, image_path: str):
    fill = slide.background.fill
    if hasattr(fill, "user_picture"):
        fill.user_picture(image_path)
        return
    # Fallback for older python-pptx versions without user_picture on FillFormat
    pic = slide.shapes.add_picture(
        image_path,
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    # SAFELY send picture to back: spTree must keep nvGrpSpPr and grpSpPr first.
    spTree = slide.shapes._spTree
    pic_el = pic._element
    # Remove current node
    spTree.remove(pic_el)
    # Find the insertion index AFTER the header nodes (usually first two)
    insert_idx = 0
    for i, child in enumerate(list(spTree)):
        tag = child.tag.rsplit('}', 1)[-1]
        if tag in ("nvGrpSpPr", "grpSpPr"):
            insert_idx = i + 1
            continue
        # After we’ve passed the header nodes, stop at the first real shape
        if insert_idx:
            break
    # Insert at the computed index (as the first actual shape)
    spTree.insert(insert_idx, pic_el)


def add_logo(slide, prs, logo_path: str, width_cm: float = 5.65):
    """Add a logo image to the lower-right corner of the slide."""
    from pptx.util import Cm
    if not os.path.exists(logo_path):
        print(f"Warning: logo file not found: {logo_path}")
        return
    logo_width = Cm(width_cm)
    logo_height = None  # auto-scale
    # Add first to get actual rendered size (height auto-scales from width)
    pic = slide.shapes.add_picture(logo_path, Cm(0), Cm(0), width=logo_width, height=logo_height)

    # Reposition to lower-right with margins; ensures it stays on-screen
    horizontal_margin = Cm(1.0)
    vertical_margin=Cm(4.0)
    print(f"DEBUG: {prs.slide_width=}, {pic.width=}, {horizontal_margin=}")
    print(f"Debug: {prs.slide_height=}, {pic.height=}, {vertical_margin=}")
    pic.left = prs.slide_width - pic.width - horizontal_margin
    pic.top = prs.slide_height - vertical_margin


def create_title_slide(prs, title, subtitle,logo_path):

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    #print(f"DEBUG: Checking file exists: assets/backgrounds/title_background_pro.jpg: {os.path.exists('assets/backgrounds/title_background_pro.jpg')}")
    #print(f"DEBUG: WHere am I: {os.getcwd()}")

    IMG = pathlib.Path(__file__).parents[1] / "assets" / "backgrounds" / "title_background_pro.png"
    set_slide_background_picture(slide, prs, str(IMG))
    #set_slide_background_picture(slide, prs, "assets/backgrounds/title_background_pro.png")

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(60)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(36)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_logo(slide, prs, logo_path )
    
def create_content_side(prs, slide_data,logo_path):
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = slide_data.get('title', 'Slide')
    # Set blue background for title
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0, 51, 153)  # dark blue
    # Set title text color to white
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)

    tf = body_shape.text_frame
    tf.clear()  # Clear any existing content and start fresh
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    first_para_used = False
    for block in slide_data.get('content_blocks', []):
        btype = str(block.get('type', '')).strip().lower()
        text = block.get('body', '')

        if btype == 'text':
            for line in text.splitlines():
                if line.strip() == '':
                    continue  # skip empty lines
                p = tf.paragraphs[0] if not first_para_used else tf.add_paragraph()
                p.text = line
                p.level = 0  # ensure bulleted at base level
                for run in p.runs:
                    run.font.size = Pt(24)
                first_para_used = True

        elif btype == 'code':
            p = tf.paragraphs[0] if not first_para_used else tf.add_paragraph()
            p.text = text
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(24)
                run.font.name = 'Courier New'
                run.font.color.rgb = RGBColor(0, 0, 255)
            first_para_used = True

        elif btype == 'image':
            # Placeholder: image handling not implemented yet
            print(f"TODO: Image block not implemented yet. query: {block.get('query')} ")

    add_logo(slide, prs, logo_path)
    
    
def create_one_presentation(slide_json, theme, output_fname):
    prs = Presentation()
    print(f"DEBUG DE:\n\n{slide_json=}\n\n")
    logo_path = str(pathlib.Path(__file__).parents[1] / "assets" / "logos" / "logoPro.png")
    
    create_title_slide(prs, slide_json['title'], slide_json.get('subtitle', ''),logo_path)
    for slide_data in slide_json['slides']:
        create_content_side(prs, slide_data,logo_path)
    output_path = os.path.join("output", output_fname)
    prs.save(output_path)

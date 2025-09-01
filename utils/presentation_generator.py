"""Module for generating PowerPoint presentations using OpenAI and python-pptx."""

import os
import re
from pathlib import Path
from openai import OpenAI
import json
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pygments import lex
from pygments.lexers import get_lexer_by_name
from pygments.styles import get_style_by_name
import traceback
from io import BytesIO
import requests

# Constants
ASSETS_DIR = Path(__file__).parent.parent / 'assets'
TEMPLATE_DIR = ASSETS_DIR / 'templates'
LOGOS_DIR = ASSETS_DIR / 'logos'

THEMES = {
    'theme1': {
        'name': 'Elementary',
        'template': 'elementary.pptx',
        'font': 'Comic Sans MS',
        'code_font': 'Courier New'
    },
    'theme2': {
        'name': 'Middle',
        'template': 'middle.pptx',
        'font': 'Arial',
        'code_font': 'Courier New'
    },
    'theme3': {
        'name': 'Professional',
        'template': 'professional.pptx',
        'font': 'Calibri',
        'code_font': 'Consolas'
    },
    'theme4': {
        'name': 'WhiteLabel',
        'template': 'whitelabel.pptx',
        'font': 'Calibri',
        'code_font': 'Consolas'
    }
}

class PresentationError(Exception):
    """Custom exception for presentation generation errors."""
    pass


def _fetch_image(block: dict) -> BytesIO | None:
    """Retrieve an image for an image content block.

    The block may specify either a direct ``url`` or a ``query`` to search for a
    relevant image. When a query is provided, the Unsplash source endpoint is
    used to fetch a representative image without requiring an API key.

    Parameters
    ----------
    block:
        Content block containing image information.

    Returns
    -------
    BytesIO | None
        Image data as a stream if retrieval succeeds, otherwise ``None``.
    """
    url = block.get("url")
    if not url and (query := block.get("query")):
        url = f"https://source.unsplash.com/1600x900/?{requests.utils.quote(query)}"
    if not url:
        return None
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return BytesIO(response.content)
    except Exception:
        return None

def _parse_slide_deck_prompt(prompt: str) -> dict | None:
    """Parse a user prompt that already describes a slide deck.

    The supported format is a simple text outline for a presentation:

    ```
    Title: My Presentation
    Subtitle: Optional subtitle
    Slide 1: Introduction
    - First bullet
    - Second bullet
    Slide 2: Next Topic
    * Another point
    ```

    Parameters
    ----------
    prompt:
        Raw user input that may describe a slide deck.

    Returns
    -------
    dict | None
        Parsed content matching the structure expected by the rest of the
        application, or ``None`` if the prompt does not match the expected
        outline format.
    """
    lines = [line.rstrip() for line in prompt.strip().splitlines() if line.strip()]

    title = None
    subtitle = None
    slides = []
    current_title = None
    bullets: list[str] = []

    for line in lines:
        if match := re.match(r"^Title:\s*(.+)", line, re.IGNORECASE):
            title = match.group(1).strip()
            continue
        if match := re.match(r"^Subtitle:\s*(.+)", line, re.IGNORECASE):
            subtitle = match.group(1).strip()
            continue
        if match := re.match(r"^Slide\s*\d+:\s*(.+)", line, re.IGNORECASE):
            if current_title and bullets:
                slides.append(
                    {
                        "id": f"slide{len(slides)+1}",
                        "title": current_title,
                        "content_blocks": [
                            {"type": "text", "body": b} for b in bullets
                        ],
                    }
                )
            current_title = match.group(1).strip()
            bullets = []
            continue
        if current_title and (
            match := re.match(r"^(?:[-*]|\d+[.)])\s+(.*)", line)
        ):
            bullets.append(match.group(1).strip())

    if current_title and bullets:
        slides.append(
            {
                "id": f"slide{len(slides)+1}",
                "title": current_title,
                "content_blocks": [
                    {"type": "text", "body": b} for b in bullets
                ],
            }
        )

    if title and slides:
        content = {"title": title, "slides": slides}
        if subtitle:
            content["subtitle"] = subtitle
        return content
    return None


def _code_font_size(code: str) -> Pt:
    """Determine an appropriate font size for a code snippet."""
    lines = code.splitlines()
    max_len = max((len(line) for line in lines), default=0)
    size = 20
    if len(lines) > 10 or max_len > 60:
        size = 16
    if len(lines) > 15 or max_len > 80:
        size = 14
    return Pt(size)


def _add_code_paragraph(tf, code: str, language: str, theme: dict) -> None:
    """Add a syntax highlighted code block to a text frame."""
    lexer = get_lexer_by_name(language or "python")
    style = get_style_by_name("default")
    font_size = _code_font_size(code)

    p = tf.add_paragraph()
    p._element.get_or_add_pPr().insert(0, OxmlElement("a:buNone"))

    for ttype, value in lex(code, lexer):
        run = p.add_run()
        run.text = value
        run.font.name = theme["code_font"]
        run.font.size = font_size
        color = style.style_for_token(ttype)["color"]
        if color:
            r, g, b = (int(color[i:i+2], 16) for i in (0, 2, 4))
            run.font.color.rgb = RGBColor(r, g, b)

def create_one_presentation(content: dict, theme_name: str, output_file: str):
    """Create a PowerPoint presentation from the generated content."""
    if theme_name not in THEMES:
        raise ValueError(f"Theme '{theme_name}' not found")
    
    theme = THEMES[theme_name]
    template_file = TEMPLATE_DIR / theme['template']
    
    if not template_file.exists():
        raise FileNotFoundError(f"Template file not found: {template_file}")
    
    try:
        prs = Presentation(template_file)
        
        # Get title slide layout
        title_layout = None
        content_layout = None
        
        for layout in prs.slide_layouts:
            if layout.name == "TITLE":
                title_layout = layout
            elif layout.name == "TITLE_AND_BODY":
                content_layout = layout
        
        if not title_layout or not content_layout:
            raise PresentationError("Required slide layouts not found in template")
        
        # Create title slide with custom formatting
        title_slide = prs.slides.add_slide(title_layout)
        
        # Add title textbox with enhanced styling
        left = int(prs.slide_width * 0.1)  # 10% from left
        top = int(prs.slide_height * 0.4)   # 40% from top
        width = int(prs.slide_width * 0.8)  # 80% of slide width
        height = int(prs.slide_height * 0.2)  # 20% of slide height
        
        title_box = title_slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        
        title_para = title_frame.add_paragraph()
        title_para.text = content['title']
        title_para.font.size = Pt(44)
        title_para.font.name = theme['font']
        title_para.alignment = 1  # Center align
        
        # Add subtitle
        if 'subtitle' in content:
            top = int(prs.slide_height * 0.7)  # 70% from top
            height = int(prs.slide_height * 0.1)  # 10% of slide height
            
            subtitle_box = title_slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            subtitle_frame.word_wrap = True
            
            subtitle_para = subtitle_frame.add_paragraph()
            subtitle_para.text = content['subtitle']
            subtitle_para.font.size = Pt(32)
            subtitle_para.font.name = theme['font']
            subtitle_para.alignment = 1  # Center align
        
        # Create content slides
        for slide_data in content['slides']:
            slide = prs.slides.add_slide(content_layout)
            
            # Add title
            title_shape = None
            body_shape = None
            
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Title
                    title_shape = shape
                elif shape.placeholder_format.type == 2:  # Body
                    body_shape = shape
            
            if title_shape:
                title_shape.text = slide_data['title']
            
            if body_shape:
                has_image = any(block.get('type') == 'image' for block in slide_data['content_blocks'])
                if has_image:
                    body_shape.width = int(prs.slide_width * 0.6)

                tf = body_shape.text_frame
                tf.clear()

                image_top = int(prs.slide_height * 0.25)

                for block in slide_data['content_blocks']:
                    if block['type'] == 'text':
                        p = tf.add_paragraph()
                        p.text = block['body']
                        p.font.name = theme['font']
                        p.font.size = Pt(24)
                    elif block['type'] == 'code':
                        _add_code_paragraph(
                            tf, block['body'], block.get('language', 'python'), theme
                        )
                    elif block['type'] == 'image':
                        img_stream = _fetch_image(block)
                        if img_stream:
                            pic = slide.shapes.add_picture(
                                img_stream,
                                int(prs.slide_width * 0.65),
                                image_top,
                                width=int(prs.slide_width * 0.3),
                            )
                            image_top = pic.top + pic.height + Inches(0.2)

        
        # Save the presentation
        prs.save(output_file)
        
    except Exception as e:
        raise PresentationError(f"Error creating presentation: {str(e)}")

def generate_teacher_guide(content: dict) -> str:
    """Generate a teacher guide from the presentation content."""
    guide = f"# Teacher Guide: {content['title']}\n\n"
    
    if 'subtitle' in content:
        guide += f"## {content['subtitle']}\n\n"
    
    guide += "## Slide-by-Slide Notes\n\n"
    
    for i, slide in enumerate(content['slides'], 1):
        guide += f"### Slide {i}: {slide['title']}\n\n"
        
        for block in slide['content_blocks']:
            if block['type'] == 'text':
                guide += f"- {block['body']}\n"
            elif block['type'] == 'code':
                guide += f"\nCode example ({block.get('language', 'code')}):\n```{block.get('language', '')}\n{block['body']}\n```\n"
            elif block['type'] == 'image':
                desc = block.get('caption') or block.get('query') or block.get('url', '')
                guide += f"Image: {desc}\n"
        
        guide += "\n"
    
    return guide

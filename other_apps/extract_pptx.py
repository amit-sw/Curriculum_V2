import zipfile, os, pathlib
from pptx import Presentation
import json
from lxml import etree

def list_layouts_in_order(pptx_path):
    prs = Presentation(pptx_path)

    # Prefer all masters (PowerPoint files can have multiple)
    masters = getattr(prs, "slide_masters", None)
    if masters is None:
        # Older python-pptx: fall back to the single default master
        masters = [prs.slide_master]

    result = []
    for m_idx, master in enumerate(masters):
        master_name = getattr(master, "name", None) or f"Master {m_idx+1}"
        for l_idx, layout in enumerate(master.slide_layouts):
            layout_name = (layout.name or f"Layout {l_idx+1}").strip()
            result.append({
                "master_index": m_idx,
                "layout_index": l_idx,
                "master_name": master_name,
                "layout_name": layout_name,
            })
    return result

def extract_all_images_fast(pptx_path, out_dir="extracted_images"):
    out = pathlib.Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(pptx_path) as z:
        for m in z.namelist():
            if m.startswith("ppt/media/"):
                z.extract(m, out)
    print(f"Done. Check {out / 'ppt' / 'media'}")

def extract_shapes_info(pptx_path, out_dir="extracted_shapes"):
    out = pathlib.Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    prs = Presentation(pptx_path)
    shapes_data = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            shape_info = {
                "slide_number": i,
                "shape_type": shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
                "name": shape.name,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "text": shape.text if hasattr(shape, "text") else None,
            }
            if shape.shape_type.name == "PICTURE":
                image = shape.image
                shape_info["image"] = {
                    "filename": image.filename if hasattr(image, "filename") else None,
                    "content_type": image.content_type if hasattr(image, "content_type") else None,
                    "width": image.width if hasattr(image, "width") else None,
                    "height": image.height if hasattr(image, "height") else None,
                }
            shapes_data.append(shape_info)
    json_path = out / (pathlib.Path(pptx_path).stem + ".json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(shapes_data, f, indent=2, ensure_ascii=False)
    print(f"Shapes info extracted to {json_path}")

def extract_shapes_info_xml(pptx_path, out_dir="extracted_shapes_xml"):
    """
    Extracts shape information from slide XMLs in a pptx, using lxml and XPath.
    For each shape (p:sp or p:pic), extracts slide number, name, text, tag, and for pictures, the r:embed id.
    Writes all shapes' info to a JSON file in out_dir.
    """
    out = pathlib.Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    pptx_stem = pathlib.Path(pptx_path).stem
    shapes_data = []
    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    }
    with zipfile.ZipFile(pptx_path) as z:
        slide_files = sorted([n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")])
        for idx, slide_file in enumerate(slide_files, start=1):
            xml_bytes = z.read(slide_file)
            root = etree.fromstring(xml_bytes)
            # Find all shape elements (p:sp) and picture elements (p:pic)
            for el in root.xpath(".//p:sp | .//p:pic", namespaces=nsmap):
                tag_local = etree.QName(el.tag).localname
                # Get name from p:nvSpPr/p:cNvPr or p:nvPicPr/p:cNvPr
                cNvPr = el.find(".//p:cNvPr", namespaces=nsmap)
                name = cNvPr.get("name") if cNvPr is not None else None
                # Extract text by concatenating all a:t descendants
                text_nodes = el.xpath(".//a:t", namespaces=nsmap)
                text = "".join(t.text for t in text_nodes if t.text)
                shape_info = {
                    "slide_number": idx,
                    "name": name,
                    "text": text,
                    "shape_tag": tag_local,
                }
                if tag_local == "pic":
                    # For pictures, get r:embed from a:blip
                    blip = el.find(".//a:blip", namespaces=nsmap)
                    embed = blip.get("{%s}embed" % nsmap["r"]) if blip is not None else None
                    shape_info["r_embed"] = embed
                shapes_data.append(shape_info)
    json_path = out / f"{pptx_stem}.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(shapes_data, f, indent=2, ensure_ascii=False)
    print(f"XML shapes info extracted to {json_path}")

#extract_all_images_fast("assets/templates/professional.pptx")
#extract_all_images_fast("assets/templates/elementary.pptx")
#extract_all_images_fast("assets/templates/middle.pptx")
#extract_all_images_fast("assets/templates/whitelabel.pptx")

#extract_shapes_info("assets/templates/elementary.pptx")
#extract_shapes_info("assets/templates/middle.pptx")
#extract_shapes_info("assets/templates/whitelabel.pptx")

#extract_shapes_info_xml("assets/templates/elementary.pptx")
#extract_shapes_info_xml("assets/templates/middle.pptx")
#extract_shapes_info_xml("assets/templates/whitelabel.pptx")

list_layouts = list_layouts_in_order("assets/templates/elementary.pptx")
print("Layouts in Elementary template:")
for l in list_layouts:
    print(l)
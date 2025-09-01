import sys, os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pathlib

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm
import requests
from io import BytesIO

#from utils.presentation_generator import create_one_presentation

my_presentation={
  "title": "Memory in Agentic AI",
  "subtitle": "Two-page deck: concepts, architecture, and practical patterns",
  "slides": [
    {
      "type": "Text",
      "body": "Slide 1 \u2014 Memory in Agentic AI: Concepts and Architecture\n- Agentic AI relies on multiple memory kinds to support autonomy: working memory, episodic memory, and semantic memory.\n- Working memory holds current goals, plans, and salient context during task execution.\n- Episodic memory stores past experiences, events, and action outcomes for learning and recall.\n- Semantic memory encodes policies, world models, and domain knowledge for planning and decision making.\n- Retrieval is powered by embeddings and similarity search to fetch relevant memories during planning."
    },
    {
      "type": "Image",
      "query": "diagram of memory layers in agentic AI",
      "caption": "Memory layers: working, episodic, and semantic memory arranged to support planning"
    },
    {
      "type": "Text",
      "body": "Slide 2 \u2014 Patterns and Practical Implementations\n- Memory stores: ephemeral working memory vs persistent long-term store; retrieval pipeline.\n- Memory lifecycle: decay, pruning, forgetting, consolidation to manage memory footprint.\n- Retrieval patterns: embedding-based search, context windows, episodic retrieval at planning time.\n- Safety and privacy: data minimization, access controls, auditing, redaction.\n- Minimal Python memory-store example provided below."
    },
    {
      "type": "Code",
      "language": "Python",
      "body": "import time\nimport math\n\nclass MemoryStore:\n    def __init__(self):\n        self.memories = []\n\n    def add(self, item, embedding=None, timestamp=None):\n        self.memories.append({'item': item, 'emb': embedding, 'ts': timestamp or time.time()})\n\n    def _cosine(self, a, b):\n        if a is None or b is None:\n            return 0.0\n        dot = sum(x*y for x,y in zip(a,b))\n        na = math.sqrt(sum(x*x for x in a))\n        nb = math.sqrt(sum(y*y for y in b))\n        if na == 0 or nb == 0:\n            return 0.0\n        return dot / (na * nb)\n\n    def query(self, q_embedding, top_k=3):\n        if not self.memories:\n            return []\n        scored = [(self._cosine(m['emb'], q_embedding), m['item']) for m in self.memories]\n        scored.sort(key=lambda x: x[0], reverse=True)\n        return [item for _, item in scored[:top_k]]\n\n# Example usage\nstore = MemoryStore()\nstore.add('Observed obstacle', [0.2, 0.5, 0.9])\nstore.add('Goal: reach waypoint B', [0.3, 0.4, 0.1])\n\nquery_vec = [0.25, 0.45, 0.2]\ntop = store.query(query_vec, top_k=2)\nprint(top)"
    },
    {
      "type": "Image",
      "query": "memory retrieval architecture for agentic AI",
      "caption": "Retrieval pipeline: store \u2192 index \u2192 query \u2192 apply"
    },
  ],
}

my_presentation={'title': 'GPT-OSS-20B Deck', 
                 'subtitle': '', 
                 'slides': [
                     {'type': 'Text', 'id': 'slide1', 'title': 'GPT-OSS-20B: Overview', 
                      'body': '- 20B-parameter transformer-based language model described by OpenAI as an open-source release\n- Aims to enable researchers and developers to experiment with LLMs without vendor lock-in\n- Supports a broad set of tasks: text generation, summarization, translation, Q&A, and multi-turn dialogue\n- Emphasizes transparency, reproducibility, and community contributions\n- Licensing terms and official docs vary by release; verify current terms before use'}, 
                      {'type': 'Text', 'id': 'slide2', 'title': 'Core Capabilities and Limitations', 
                       'body': '- Strong natural language understanding and generation across domains\n- Useful for coding help, technical explanations, and general QA\n- Handles multi-turn conversations with context windows\n- Limitations: biases, hallucinations, data cutoff, sensitivity to prompts\n- Safety and governance: implement monitoring, filters, and adhere to OSS guidelines'}, 
                      {'type': 'Text', 'id': 'slide3', 'title': 'Getting Started and Deployment', 
                       'body': '- Hardware and infrastructure: GPUs with adequate VRAM; plan for distributed inference if needed\n- Software stack: Python, PyTorch, transformers, and accelerators; obtain weights from official sources\n- Typical workflow: tokenize input, run model, decode with controlled sampling (temperature, top-k/top-p)\n- Best practices: task-specific evaluation, latency monitoring, possible quantization/distillation\n- OSS licensing and governance: comply with license terms, attribution, and community guidelines'}
                       ], 
                'user_message': 'Create a concise three-slide deck about GPT-OSS-20B from OpenAI, covering overview, capabilities/limitations, and getting started/deployment. Include a note that details may vary by release and to verify against official docs.'
                }

my_presentation = {'title': 'AI in FinTech Take 3', 'subtitle': 'Three-slide overview for quick reference', 
                   'slides': [
                       {'id': 'slide-1', 'title': 'What AI in FinTech Enables', 
                        'content_blocks': [
                            {'type': 'Text', 'body': 'AI transforms financial services by enabling smarter decisions, personalized experiences, and stronger risk controls.'}, 
                            {'type': 'Text', 'body': 'Key use cases include fraud detection, credit underwriting, algorithmic trading, customer service automation, and regulatory tech.'}, 
                            {'type': 'Text', 'body': 'Benefits: improved efficiency, scalability, and accuracy. Challenges: data quality, bias, governance, and model risk.'}, 
                            #{'type': 'Image', 'query': 'AI in FinTech architecture diagram', 'caption': 'High-level AI-enabled fintech architecture'}
                            ]}, 
                       {'id': 'slide-2', 'title': 'Core AI Techniques in FinTech', 
                        'content_blocks': [
                            {'type': 'Text', 'body': 'Data-driven models underpin most fintech AI use cases, including supervised learning for scoring and time-series for risk and pricing.'},
                            {'type': 'Text', 'body': 'Fraud detection and risk: anomaly detection, ensemble models, feature engineering, and real-time monitoring.'}, 
                            {'type': 'Text', 'body': 'Credit underwriting: logistic regression, gradient boosting, and explainability with SHAP.'}, 
                            {'type': 'Text', 'body': 'NLP and RPA for customer service and document processing; governance and monitoring are essential for reliability.'}, 
                            {'type': 'Code', 'language': 'python', 'body': "from sklearn.ensemble import IsolationForest\nimport numpy as np\n\n# Simple anomaly detection example\nX = np.random.randn(100, 2)\nclf = IsolationForest(contamination=0.1, random_state=42).fit(X)\nprint('anomaly scores:', clf.decision_function(X))"}
                            ]}, 
                       {'id': 'slide-3', 'title': 'Operationalizing AI in FinTech', 
                        'content_blocks': [
                            {'type': 'Text', 'body': 'Deployment: containerized services, scalable endpoints, and real-time scoring.'}, 
                            {'type': 'Text', 'body': 'Governance and risk: drift detection, versioning, audit trails, and explainability.'}, 
                            {'type': 'Text', 'body': 'Data privacy and security: encryption, access controls, and regulatory compliance (e.g., GDPR, PSD2).'}, 
                            {'type': 'Text', 'body': 'Monitoring and incident response: metrics, alerts, and post-deployment evaluation.'}, 
                            #{'type': 'Image', 'query': 'data drift monitoring diagram', 'caption': 'Model monitoring and drift detection in production'}
                            ]
                        },
                       ], 
                   'user_message': 'User requested a three-slide deck on AI in FinTech with concise titles, 3–5 bullets per slide, a code example, and supporting imagery.'}



def set_slide_background_picture(slide, prs, image_path: str):
    fill = slide.background.fill
    if hasattr(fill, "user_picture"):
        fill.user_picture(image_path)
        return
    # Fallback for older python-pptx versions without user_picture on FillFormat
    pic = slide.shapes.add_picture(
        image_path,
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    # SAFELY send picture to back: spTree must keep nvGrpSpPr and grpSpPr first.
    spTree = slide.shapes._spTree
    pic_el = pic._element
    # Remove current node
    spTree.remove(pic_el)
    # Find the insertion index AFTER the header nodes (usually first two)
    insert_idx = 0
    for i, child in enumerate(list(spTree)):
        tag = child.tag.rsplit('}', 1)[-1]
        if tag in ("nvGrpSpPr", "grpSpPr"):
            insert_idx = i + 1
            continue
        # After we’ve passed the header nodes, stop at the first real shape
        if insert_idx:
            break
    # Insert at the computed index (as the first actual shape)
    spTree.insert(insert_idx, pic_el)


def add_logo(slide, prs, logo_path: str, width_cm: float = 2.5):
    """Add a logo image to the lower-right corner of the slide."""
    from pptx.util import Cm
    if not os.path.exists(logo_path):
        print(f"Warning: logo file not found: {logo_path}")
        return
    logo_width = Cm(width_cm)
    logo_height = None  # auto-scale
    left = prs.slide_width - logo_width
    top = prs.slide_height - Cm(2.5)  # 2.5 cm up from bottom
    slide.shapes.add_picture(logo_path, left, top, width=logo_width, height=logo_height)


def create_title_slide(prs, title, subtitle):

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    #print(f"DEBUG: Checking file exists: assets/backgrounds/title_background_pro.jpg: {os.path.exists('assets/backgrounds/title_background_pro.jpg')}")
    #print(f"DEBUG: WHere am I: {os.getcwd()}")

    IMG = pathlib.Path(__file__).parents[1] / "assets" / "backgrounds" / "title_background_pro.png"
    set_slide_background_picture(slide, prs, str(IMG))
    #set_slide_background_picture(slide, prs, "assets/backgrounds/title_background_pro.png")

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    LOGO = pathlib.Path(__file__).parents[1] / "assets" / "logo.png"
    add_logo(slide, prs, str(LOGO))
    
def create_content_side(prs, slide_data):
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = slide_data.get('title', 'Slide')
    
    tf = body_shape.text_frame
    tf.clear()  # Clear any existing content
    for block in slide_data.get('content_blocks', []):
        if block['type'] == 'Text':
            p = tf.add_paragraph()
            p.text = block['body']
            p.font.size = Pt(18)
            p.space_after = Pt(10)
        elif block['type'] == 'Code':
            p = tf.add_paragraph()
            p.text = block['body']
            p.font.size = Pt(14)
            p.font.name = 'Courier New'
            p.font.color.rgb = RGBColor(0, 0, 255)
            p.space_after = Pt(10)
        elif block['type'] == 'Image':
            print(f"TODO TODO TODO: Not doing anything for image block yet. query: {block['query']}")

    LOGO = pathlib.Path(__file__).parents[1] / "assets" / "logo.png"
    add_logo(slide, prs, str(LOGO))
    
    
def create_one_presentation(presentation, theme, output_path):


    prs = Presentation()
    create_title_slide(prs, presentation['title'], presentation.get('subtitle', ''))
    for slide_data in presentation['slides']:
        create_content_side(prs, slide_data)
    prs.save(output_path)


create_one_presentation(my_presentation, 'theme3', 'output/Memory_in_Agentic_AI4.pptx')
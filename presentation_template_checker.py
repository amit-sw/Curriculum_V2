from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

TYPE_NAME = {v: k for k, v in PP_PLACEHOLDER.__members__.items()}

def audit_template(path):
    prs = Presentation(path)
    for i, layout in enumerate(prs.slide_layouts):
        print(f"[{i}] Layout name: {layout.name!r}")
        # List placeholders on this layout
        for ph in layout.placeholders:
            pht = getattr(ph.placeholder_format, "type", None)
            type_name = TYPE_NAME.get(pht, str(pht))
            print(f"    - ph idx={ph.placeholder_format.idx}, type={type_name}")
        print()

audit_template("assets/templates/whitelabel.pptx")